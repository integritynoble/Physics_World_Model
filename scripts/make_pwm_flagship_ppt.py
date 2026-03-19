#!/usr/bin/env python3
"""Generate PowerPoint presentation for PWM Flagship paper.

Title: Eleven Primitives and Three Gates: The Universal Structure of Computational Imaging
Target: 40-minute talk for Prof. Xin Yuan's group at Westlake University

Core message: 11 primitives + 3 gates = universal grammar for
  (1) Validating existing imaging systems
  (2) Correcting existing imaging systems
  (3) Designing NEW imaging systems
"""

import sys
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "papers" / "pwm_flagship"
FIG = OUT / "figures"

# ── Colors ────────────────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BLUE    = RGBColor(0x2E, 0x5F, 0xA1)
ORANGE  = RGBColor(0xE8, 0x71, 0x2B)
GREEN   = RGBColor(0x3D, 0x9E, 0x3F)
RED     = RGBColor(0xC9, 0x3C, 0x3C)
PURPLE  = RGBColor(0x7B, 0x4E, 0xA3)
GRAY    = RGBColor(0x6B, 0x6B, 0x6B)
GATE1   = RGBColor(0x4A, 0x90, 0xD9)
GATE2   = RGBColor(0xE8, 0xA8, 0x38)
GATE3   = RGBColor(0xD9, 0x4A, 0x4A)
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_SLIDE = RGBColor(0xF5, 0xF7, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 32
sn = 0

# ── Helpers ───────────────────────────────────────────────────────
def add_bg(slide, color=BG_SLIDE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, fn='Calibri'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = fn; p.alignment = align
    return tf

def bullets(slide, l, t, w, h, items, sz=18, color=BLACK, sp=Pt(6)):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = 'Calibri'; p.space_after = sp
    return tf

def title_bar(slide, title, subtitle=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    tb(slide, 0.6, 0.15, 12, 0.7, title, sz=32, bold=True, color=WHITE)
    if subtitle:
        tb(slide, 0.6, 0.7, 12, 0.4, subtitle, sz=16, color=RGBColor(0xCC, 0xDD, 0xFF))

def section_slide(slide, title, part=None):
    add_bg(slide, DARK)
    if part:
        tb(slide, 0.6, 2.0, 12, 0.8, part, sz=24, color=GATE1)
    tb(slide, 0.6, 2.7, 12, 1.5, title, sz=48, bold=True, color=WHITE)

def snum(slide):
    tb(slide, 12.3, 7.0, 1.0, 0.4, f"{sn}/{TOTAL}", sz=12, color=GRAY, align=PP_ALIGN.RIGHT)

def box(slide, l, t, w, h, fill_rgb, line_rgb):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_rgb; s.line.color.rgb = line_rgb
    return s

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def pic(slide, name, l, t, w, h):
    p = FIG / name
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))

# ══════════════════════════════════════════════════════════════════
# 1: TITLE
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s, BG_DARK)
tb(s, 0.8, 1.0, 11.7, 1.5, "Eleven Primitives and Three Gates:", sz=44, bold=True, color=WHITE)
tb(s, 0.8, 2.0, 11.7, 1.0, "The Universal Structure of Computational Imaging", sz=40, bold=True, color=GATE1)
tb(s, 0.8, 3.5, 11.7, 0.5, "Chengshuai Yang, David J. Brady, Xin Yuan", sz=24, color=RGBColor(0xBB,0xBB,0xBB))
tb(s, 0.8, 4.2, 11.7, 0.8, "NextGen PlatformAI  |  University of Arizona  |  Westlake University", sz=18, color=GRAY)
# Three pillars
box(s, 1.5, 5.3, 3.0, 0.9, RGBColor(0x1E,0x3A,0x5F), GATE1)
tb(s, 1.5, 5.35, 3.0, 0.8, "DIAGNOSE", sz=26, bold=True, color=GATE1, align=PP_ALIGN.CENTER)
box(s, 5.2, 5.3, 3.0, 0.9, RGBColor(0x3A,0x1E,0x1E), GATE3)
tb(s, 5.2, 5.35, 3.0, 0.8, "CORRECT", sz=26, bold=True, color=GATE3, align=PP_ALIGN.CENTER)
box(s, 8.9, 5.3, 3.0, 0.9, RGBColor(0x1E,0x3A,0x1E), GREEN)
tb(s, 8.9, 5.35, 3.0, 0.8, "DESIGN NEW", sz=26, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tb(s, 0.8, 6.5, 11.7, 0.5, "Talk for Prof. Yuan's Group  |  March 2026", sz=18, color=GRAY)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 2: OUTLINE
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Talk Outline", "~40 minutes")
bullets(s, 0.8, 1.4, 11.7, 5.5, [
    "1.  Motivation: Why do real systems fail? How to design new ones? (5 min)",
    "2.  Our Answer: 11 Primitives + 3 Gates (5 min)",
    "3.  The Finite Primitive Basis (Theorem 1) (7 min)",
    "      OperatorGraph, 11 primitives, DAG examples, basis saturation",
    "4.  The Triad Decomposition (Theorem 2) (5 min)",
    "      Three gates, lifecycle prediction, four-scenario protocol",
    "5.  Application 1: DIAGNOSE existing systems (3 min)",
    "6.  Application 2: CORRECT existing systems (5 min)",
    "      CASSI deep dive, cross-carrier results (+0.8 to +13.9 dB)",
    "7.  Application 3: DESIGN NEW systems (7 min)",
    "      Gates as design specs, primitive composition, design exploration",
    "8.  Discussion and Future Directions (3 min)",
], sz=19, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 3: THE PROBLEM
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Motivation: Two Open Problems in Computational Imaging")

# Left column: Why do real systems fail?
box(s, 0.4, 1.3, 6.0, 5.0, RGBColor(0xFA,0xE8,0xE8), GATE3)
tb(s, 0.6, 1.35, 5.6, 0.5, "Problem 1: Why Do Real Systems Fail?", sz=20, bold=True, color=GATE3)
bullets(s, 0.6, 1.85, 5.6, 4.3, [
    "Algorithms keep advancing: CS -> PnP -> Unrolling -> Transformers",
    "Yet they routinely fail on deployed systems.",
    "",
    "Root cause: forward model mismatch (H_nom != H_true)",
    "Solver faithfully solves the WRONG problem.",
    "",
    "Failures are silent, systematic, universal.",
    "Per-modality calibration exists but NO universal framework.",
], sz=15, color=DARK, sp=Pt(4))

# Right column: How to design new systems?
box(s, 6.9, 1.3, 6.0, 5.0, RGBColor(0xE0,0xFA,0xE0), GREEN)
tb(s, 7.1, 1.35, 5.6, 0.5, "Problem 2: How to Design New Systems?", sz=20, bold=True, color=GREEN)
bullets(s, 7.1, 1.85, 5.6, 4.3, [
    "Designing a new system is an open-ended, intuition-driven search.",
    "No common language across modalities (CT, MRI, CASSI...).",
    "",
    "Key unknowns at design time:",
    "  - How many measurements are enough?",
    "  - What photon budget is needed?",
    "  - How precise must calibration be?",
    "No way to evaluate BEFORE hardware fabrication.",
], sz=15, color=DARK, sp=Pt(4))

# Bottom: Our answer
box(s, 0.4, 6.5, 12.5, 0.6, RGBColor(0xE0,0xEE,0xFA), BLUE)
tb(s, 0.6, 6.55, 12.1, 0.5, "Our answer: 11 Primitives + 3 Gates = universal grammar to DIAGNOSE, CORRECT, and DESIGN NEW", sz=19, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 4: DIVERSE SYSTEMS
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Diverse Systems, Hidden Universality")
bullets(s, 0.8, 1.4, 5.8, 5.5, [
    "5 carrier families:",
    "  - Optical photons (CASSI, CACTI, SPC...)",
    "  - X-ray photons (CT, CBCT)",
    "  - Electrons (ptychography, cryo-EM)",
    "  - Nuclear spins (MRI)",
    "  - Acoustic waves (ultrasound)",
    "",
    "All share the same pipeline:",
    "  source -> propagation -> interaction",
    "       -> encoding -> detection",
    "",
    "This universality is NOT accidental.",
    "It reflects the finite structure of",
    "physical measurement processes.",
], sz=18, color=DARK)
bullets(s, 7.0, 1.4, 5.8, 5.5, [
    "Three fundamental questions:",
    "",
    "Q1: What are the minimal building",
    "    blocks of imaging physics?",
    "    --> 11 Primitives",
    "",
    "Q2: Why do reconstructions fail?",
    "    --> 3 Gates (exhaustive diagnosis)",
    "",
    "Q3: Can we use this to design",
    "    NEW imaging systems?",
    "    --> YES: compositional design",
    "    from primitives + gate constraints",
], sz=18, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 5: BIG PICTURE - 3 PILLARS
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Our Answer: A Universal Grammar for Computational Imaging")

# Three pillar boxes
# DIAGNOSE
box(s, 0.5, 1.4, 3.8, 3.5, RGBColor(0xE0,0xEE,0xFA), GATE1)
tb(s, 0.7, 1.5, 3.4, 0.4, "DIAGNOSE", sz=24, bold=True, color=GATE1, align=PP_ALIGN.CENTER)
bullets(s, 0.7, 2.0, 3.4, 2.7, [
    "Existing Instruments",
    "",
    "Triad diagnoses the",
    "dominant failure gate:",
    " - Gate 1: not enough info?",
    " - Gate 2: too much noise?",
    " - Gate 3: wrong model?",
    "",
    "Tells you WHERE the dB",
    "are hiding.",
], sz=14, color=DARK, sp=Pt(2))

# CORRECT
box(s, 4.8, 1.4, 3.8, 3.5, RGBColor(0xFA,0xE0,0xE0), GATE3)
tb(s, 5.0, 1.5, 3.4, 0.4, "CORRECT", sz=24, bold=True, color=GATE3, align=PP_ALIGN.CENTER)
bullets(s, 5.0, 2.0, 3.4, 2.7, [
    "Existing Instruments",
    "",
    "Solver-agnostic correction:",
    "fix the operator, not solver.",
    "",
    "Any solver benefits without",
    "retraining:",
    " GAP-TV, MST-L, HDNet...",
    "",
    "+0.8 to +13.9 dB recovery",
    "across 5 carrier families.",
], sz=14, color=DARK, sp=Pt(2))

# DESIGN NEW
box(s, 9.1, 1.4, 3.8, 3.5, RGBColor(0xE0,0xFA,0xE0), GREEN)
tb(s, 9.3, 1.5, 3.4, 0.4, "DESIGN NEW", sz=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
bullets(s, 9.3, 2.0, 3.4, 2.7, [
    "New Instruments",
    "",
    "Compositional design language:",
    "select primitives, connect DAG,",
    "evaluate gates BEFORE building.",
    "",
    "Gate 1: min sampling density",
    "Gate 2: min carrier budget",
    "Gate 3: max calibration error",
    "",
    "From task to hardware specs.",
], sz=14, color=DARK, sp=Pt(2))

# Foundation
box(s, 0.5, 5.2, 12.4, 1.8, RGBColor(0xF0,0xF0,0xF8), BLUE)
tb(s, 0.8, 5.3, 5.5, 0.4, "Theorem 1: Finite Primitive Basis", sz=20, bold=True, color=BLUE)
bullets(s, 0.8, 5.7, 5.5, 1.2, [
    "Every forward model = DAG over 11 primitives",
    "Sufficient, minimal, covers 170 modalities",
], sz=16, color=DARK, sp=Pt(2))
tb(s, 7.0, 5.3, 5.5, 0.4, "Theorem 2: Triad Decomposition", sz=20, bold=True, color=ORANGE)
bullets(s, 7.0, 5.7, 5.5, 1.2, [
    "Every failure = G1 + G2 + G3 (exhaustive)",
    "Gate 3 dominant in deployed systems (14/14)",
], sz=16, color=DARK, sp=Pt(2))
snum(s)

# ══════════════════════════════════════════════════════════════════
# 6: Figure 1
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Figure 1: Universal Grammar --- Validate, Correct, and Design New")
pic(s, "fig1_overview.png", 0.5, 1.3, 12.3, 5.8)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 7: SECTION - Finite Primitive Basis
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide()
section_slide(s, "The Finite Primitive Basis", "THEOREM 1")
tb(s, 0.6, 4.0, 12, 1.0, "Every imaging forward model decomposes into a DAG over exactly 11 primitives", sz=22, color=GRAY)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 8: OperatorGraph
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "The OperatorGraph Intermediate Representation")
bullets(s, 0.8, 1.4, 11.7, 5.5, [
    "Every imaging forward model encoded as a typed DAG:",
    "  - Each node wraps a single primitive physical operator",
    "  - Edges define data flow from source to detector",
    "  - Each edge carries tensor shape + dtype metadata (static validation)",
    "",
    "Every linear primitive implements:",
    "  - forward(x) -> y    (physics simulation)",
    "  - adjoint(y) -> x    (validated: <Hx,y> = <x,H*y> to delta < 10^-6)",
    "",
    "Mathematical safety rails (NOT a neural network):",
    "  - Adjoint consistency: physical energy accounting certificate",
    "  - Lipschitz bounds: nonlinear primitives restricted to 5 canonical families",
    "  - Not a universal approximator --> falsifiable, interpretable, certifiable",
    "",
    "If a new modality can't be represented --> extension protocol fires",
    "(This is how Scatter R was discovered when Compton imaging gave e=0.34)",
], sz=19, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 9: 11 Primitives
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "The 11 Canonical Primitives")
prims = [
    "  #    Primitive       Notation            Physical Action",
    "  1    Propagate       P(d, lambda)         Free-space wave propagation",
    "  2    Modulate        M(m)                Element-wise multiplication (mask, coil)",
    "  3    Project         Pi(theta)            Radon line-integral projection",
    "  4    Encode          F(k)                Fourier-domain encoding (k-space)",
    "  5    Convolve        C(h)                Spatial convolution (PSF)",
    "  6    Accumulate      Sigma               Sum over spectral/temporal axis",
    "  7    Detect          D(g, eta)            Detector response (5 families)",
    "  8    Sample          S(Omega)            Sub-sampling on index set",
    "  9    Disperse        W(alpha, a)          Wavelength-dependent shift",
    " 10    Scatter         R(sigma, dE)         Direction change / energy shift",
    " 11    Transform       Lambda(f, theta)     Pointwise nonlinear physics",
]
tb(s, 0.5, 1.3, 12.3, 0.5, prims[0], sz=16, bold=True, color=BLUE, fn='Consolas')
bullets(s, 0.5, 1.8, 12.3, 4.0, prims[1:], sz=15, color=DARK, sp=Pt(2))
# Key points
bullets(s, 0.5, 5.4, 12.3, 2.0, [
    "6 physics-stage families: propagation, elastic interaction, inelastic interaction (scatter),",
    "   pointwise nonlinear, encoding-projection, detection-readout",
    "Minimal basis: removing any primitive leaves >= 1 modality with error > epsilon",
    "Covers 170 modalities across 5 carrier families + particle beams (neutron, proton, muon CT)",
], sz=15, color=GRAY, sp=Pt(2))
snum(s)

# ══════════════════════════════════════════════════════════════════
# 10: Figure 2
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Figure 2: OperatorGraph DAGs, Fidelity Ladder, Basis Growth")
pic(s, "fig2_operatorgraph.png", 0.5, 1.3, 12.3, 5.8)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 11: DAG Examples
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "DAG Decompositions: Same Language, Every Carrier")
bullets(s, 0.8, 1.3, 6.0, 5.5, [
    "Optical photons:",
    "  CASSI:  M -> W -> Sigma -> D",
    "  CACTI:  M -> Sigma -> D",
    "  SPC:    M -> Sigma -> D",
    "  Lensless: C -> D",
    "",
    "X-ray photons:",
    "  CT:     Pi -> D",
    "  CBCT:   Pi -> D",
    "",
    "Electrons:",
    "  Ptycho: M -> P -> D",
    "  Cryo:   P -> Lambda -> D",
], sz=18, color=DARK)
bullets(s, 7.0, 1.3, 5.8, 5.5, [
    "Nuclear spins:",
    "  MRI:    M -> F -> S -> D",
    "",
    "Acoustic waves:",
    "  US:     P -> Sigma -> D",
    "",
    "Exotic / held-out (frozen library):",
    "  Ghost imaging: M -> Sigma -> D",
    "  THz-TDS: C -> D_coh",
    "  Compton: M -> R -> D",
    "  DOT:    M -> R o P o R -> D",
    "",
    "Median: 3 nodes, depth 3",
    "All e_img < 0.01 (1% relative error)",
], sz=18, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 12: Basis saturation
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Basis-Growth Saturation: K=11 Suffices for N=170")
pic(s, "fig9_basis_growth.png", 0.5, 1.3, 6.0, 5.5)
bullets(s, 7.0, 1.4, 5.5, 5.5, [
    "K=11 at N=35 modalities",
    "No new primitive needed for",
    "next 135 modalities!",
    "",
    "Including particle-beam probes:",
    "  Neutron CT, Proton CT,",
    "  Muon tomography",
    "  --> 0 new primitives",
    "",
    "Consistent with theorem:",
    "  Once all 6 physics-stage",
    "  families are covered,",
    "  basis is complete.",
    "",
    "New modality = new combination",
    "of existing primitives.",
], sz=18, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 13: SECTION - Triad
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide()
section_slide(s, "The Triad Decomposition", "THEOREM 2")
tb(s, 0.6, 4.0, 12, 1.0, "Every reconstruction failure has exactly 3 independent root causes", sz=22, color=GRAY)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 14: Three Gates
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Three Gates: Exhaustive Failure Decomposition")

# G1
box(s, 0.5, 1.3, 3.8, 2.5, RGBColor(0xE0,0xEE,0xFA), GATE1)
tb(s, 0.7, 1.4, 3.4, 0.4, "Gate 1: Recoverability", sz=20, bold=True, color=GATE1)
bullets(s, 0.7, 1.9, 3.4, 1.8, [
    "Is enough information captured?",
    "Compression ratio gamma = m/n",
    "Null space --> lost forever",
    "Fix: more measurements",
    "Lifecycle: DESIGN stage",
], sz=14, color=DARK, sp=Pt(2))

# G2
box(s, 4.8, 1.3, 3.8, 2.5, RGBColor(0xFA,0xF0,0xE0), GATE2)
tb(s, 5.0, 1.4, 3.4, 0.4, "Gate 2: Carrier Budget", sz=20, bold=True, color=GATE2)
bullets(s, 5.0, 1.9, 3.4, 1.8, [
    "Is the SNR sufficient?",
    "Shot noise, thermal, T1/T2",
    "Noise corrupts information",
    "Fix: source/detector upgrade",
    "Lifecycle: DESIGN stage",
], sz=14, color=DARK, sp=Pt(2))

# G3
box(s, 9.1, 1.3, 3.8, 2.5, RGBColor(0xFA,0xE0,0xE0), GATE3)
tb(s, 9.3, 1.4, 3.4, 0.4, "Gate 3: Operator Mismatch", sz=20, bold=True, color=GATE3)
bullets(s, 9.3, 1.9, 3.4, 1.8, [
    "Does H_nom match H_true?",
    "Misalignment, drift, model error",
    "Solver solves wrong problem",
    "Fix: calibrate the operator",
    "Lifecycle: DEPLOYMENT stage",
], sz=14, color=DARK, sp=Pt(2))

# MSE decomposition
box(s, 0.5, 4.1, 12.4, 1.0, RGBColor(0xF8,0xF8,0xFF), BLUE)
tb(s, 0.8, 4.15, 11.8, 0.9,
   "MSE  <=  MSE^(G1) + MSE^(G2) + MSE^(G3)    (exhaustive, no residual term)", sz=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# Lifecycle insight
box(s, 0.5, 5.3, 12.4, 1.8, RGBColor(0xE8,0xFE,0xE8), GREEN)
tb(s, 0.8, 5.35, 11.8, 0.4, "Lifecycle Prediction (confirmed: 14/14 configurations)", sz=20, bold=True, color=GREEN)
bullets(s, 0.8, 5.8, 11.8, 1.2, [
    "Design stage: Gates 1 & 2 are optimized (sampling geometry, carrier selection)",
    "Deployment stage: Gate 3 (mismatch) is the residual bottleneck --- this is why algorithms fail!",
    "The Triad tells you WHERE the dB are hiding, so effort targets the largest marginal return.",
], sz=16, color=DARK, sp=Pt(2))
snum(s)

# ══════════════════════════════════════════════════════════════════
# 15: Figure 3 - Triad
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Figure 3: Triad Decomposition Validated Across All Carriers")
pic(s, "fig3_triad.png", 0.5, 1.3, 12.3, 5.8)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 16: 4-Scenario Protocol
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "The Four-Scenario Evaluation Protocol")
bullets(s, 0.8, 1.4, 11.7, 5.5, [
    "  Scenario I  (Ideal):     y = H_true * x,  reconstruct with H_true",
    "       --> Oracle upper bound. Best achievable quality.",
    "",
    "  Scenario II (Mismatch):  y = H_true * x,  reconstruct with H_nom  (H_nom != H_true)",
    "       --> Standard deployment condition. The model is WRONG.",
    "",
    "  Scenario III (Oracle):   Same y as Sc. II, reconstruct with H_true",
    "       --> Correction ceiling. What's achievable if you knew the truth.",
    "",
    "  Scenario IV (Corrected): Same y as Sc. II, reconstruct with H_hat (estimated)",
    "       --> Our contribution. Autonomous calibration.",
    "",
    "Recovery ratio:  rho = (PSNR_IV - PSNR_II) / (PSNR_I - PSNR_II)",
    "  rho = 0: calibration useless.  rho = 1: fully closes the gap.",
    "",
    "This protocol applies identically across ALL modalities and carriers.",
], sz=18, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 17: SECTION - DIAGNOSE
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide()
section_slide(s, "Application 1: DIAGNOSE", "EXISTING INSTRUMENTS")
tb(s, 0.6, 4.0, 12, 1.0, "Diagnose the dominant failure gate before attempting any fix", sz=22, color=GRAY)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 18: Validate details
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Validation: Root-Cause Diagnosis Across 12 Modalities")
bullets(s, 0.8, 1.3, 5.8, 5.5, [
    "For every imaging configuration, the",
    "Triad produces a TriadReport:",
    "",
    "  - dominant_gate (enum)",
    "  - evidence_scores (3 floats)",
    "  - confidence_interval (95% CI)",
    "  - recommended_action (string)",
    "  - parameter_sensitivities (dict)",
    "",
    "Result: Gate 3 dominant in 14/14!",
    "",
    "Falsifiable predictions confirmed:",
    "  SPC: G1 at <5% compression",
    "  MRI: G2 at noise sigma > 1%",
    "  CT: pure G3 even at extremes",
    "  US: G2-immune at 50% RF noise",
], sz=17, color=DARK)
bullets(s, 7.0, 1.3, 5.8, 5.5, [
    "Cross-over thresholds:",
    "",
    "SPC: G1<->G3 at ~5% compression",
    "  Below: more data helps most",
    "  Above: fix the model instead",
    "",
    "MRI: G2<->G3 at sigma ~1%",
    "  Sharpest transition in framework",
    "  G3 dominated -> G2 dominated",
    "  in one order of magnitude",
    "",
    "These thresholds convert",
    "'is this design adequate?'",
    "into quantitative pass/fail.",
    "",
    "8/12 predictions confirmed by",
    "simulation or hardware data.",
], sz=17, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 19: SECTION - CORRECT
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide()
section_slide(s, "Application 2: CORRECT", "EXISTING INSTRUMENTS")
tb(s, 0.6, 4.0, 12, 1.0, "Fix the operator, not the solver --- any solver benefits without retraining", sz=22, color=GRAY)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 20: Correction pipeline
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Solver-Agnostic Operator Correction")
bullets(s, 0.8, 1.4, 11.7, 5.5, [
    "Key insight: correct the forward model parameters, not the solver weights.",
    "",
    "Two-stage correction pipeline (when Gate 3 is dominant):",
    "",
    "  Algorithm 1: Hierarchical Beam Search (~300s/scene)",
    "    1D sweeps -> 3D beam search -> 2D beam search -> coordinate descent",
    "    Accuracy: +/-0.1-0.2 pixels",
    "",
    "  Algorithm 2: Joint Gradient Refinement (~3200s/scene)",
    "    Differentiable forward model + Adam optimizer + multi-start restarts",
    "    Accuracy: +/-0.05-0.1 pixels (3-5x improvement over Alg. 1)",
    "",
    "Critical property: SOLVER-AGNOSTIC",
    "  The corrected operator H_hat benefits ANY downstream solver:",
    "  TwIST, GAP-TV, DGSMP, MST-L, CST-L, HDNet, EfficientSCI...",
    "  ALL improve without retraining!",
    "",
    "Better solvers + better operators = compounding gains.",
], sz=18, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 21: Figure 4 - Correction results
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Figure 4: Correction Results Across 14 Configurations, 5 Carriers")
pic(s, "fig4_correction_bar.png", 0.3, 1.3, 6.0, 5.5)
bullets(s, 6.5, 1.3, 6.3, 6.0, [
    "Optical photons (6 modalities):",
    "  CASSI:  +6.68 dB (10 scenes, MST-L)",
    "  CACTI:  +4.37 dB (6 videos)",
    "  SPC:    +2.31 dB (11 images)",
    "  Lensless: +13.9 dB",
    "  Holography: +7.41 dB",
    "  Fluorescence: +0.80 dB",
    "",
    "X-ray: CT +3.10 dB, CBCT +5.0 dB",
    "Electrons: Ptycho +6.77, Cryo +2.4 dB",
    "Spins: MRI +2.5 dB (ESPIRiT adopted)",
    "Acoustic: Ultrasound +1.8 dB",
    "",
    "Gate 3 dominant in ALL 14/14!",
    "Recovery ratio: 0.82-1.00",
    "(82-100% of oracle ceiling)",
], sz=16, color=DARK, sp=Pt(2))
snum(s)

# ══════════════════════════════════════════════════════════════════
# 22: CASSI deep dive
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "CASSI Deep Dive: Prof. Yuan's SCI Framework")
bullets(s, 0.8, 1.4, 11.7, 5.5, [
    "CASSI DAG:  M(mask) -> W(dispersion) -> Sigma(spectral sum) -> D(detect)",
    "",
    "5-parameter mismatch: (dx, dy, theta, a1, alpha)",
    "  True: dx=0.5 px, dy=0.3 px, theta=0.1 deg, a1=2.02, alpha=0.15 deg",
    "",
    "What happens when H_nom != H_true?",
    "  A sub-pixel mask shift collapses ALL solvers to ~21 dB (from 24-35 dB)",
    "  MST-L drops by 13.98 +/- 0.62 dB --- catastrophic and silent!",
    "",
    "Correction recovers quality WITHOUT retraining any solver:",
    "  GAP-TV:  Sc.II = 25.39 -> Sc.IV = 26.15 dB  (rho = 0.85)",
    "  MST-L:   Sc.II = 21.05 -> Sc.IV = 27.73 dB  (rho = 0.46 for 5-param)",
    "",
    "Key finding for SCI community:",
    "  The solver with highest ideal PSNR (MST-L) is MOST sensitive to mismatch!",
    "  Solver quality and model fidelity are complementary, not redundant.",
    "",
    "Real-data validation: TSA scenes (5 scenes, 660x660, 28 bands) confirmed.",
], sz=17, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 23: Figure 5 - CASSI
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Figure 5: CASSI Results + Hardware Validation")
pic(s, "fig5_deepdive.png", 0.3, 1.3, 6.2, 5.5)
pic(s, "fig5_hardware.png", 6.7, 1.3, 6.2, 5.5)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 24: Visual comparison
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Figure 8: Visual Comparison Across All Scenarios")
pic(s, "fig8_visual_comparison.png", 0.3, 1.3, 12.7, 5.8)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 25: Cross-carrier
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Figure 6: Cross-Carrier Zero-Shot Generalization")
pic(s, "fig6_zeroshot.png", 1.0, 1.3, 11.0, 5.8)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 26: SECTION - DESIGN NEW
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide()
add_bg(s, BG_DARK)
tb(s, 0.6, 1.5, 12, 0.8, "APPLICATION 3", sz=24, color=GREEN)
tb(s, 0.6, 2.3, 12, 1.5, "Design NEW Imaging Systems", sz=48, bold=True, color=WHITE)
tb(s, 0.6, 4.0, 12, 1.5,
   "The 11 primitives provide a compositional design language.\n"
   "The 3 gates provide quantitative design specifications.\n"
   "Evaluate a candidate system BEFORE any hardware is fabricated.",
   sz=22, color=GRAY)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 27: Gates as design specs
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Gates as Quantitative Design Specifications")
bullets(s, 0.8, 1.4, 11.7, 5.5, [
    "The three gates, when INVERTED, yield design constraints:",
    "",
    "  Gate 1 --> Minimum measurement budget",
    "    Given target signal class + quality, what's the minimum compression ratio?",
    "    Example: SPC needs gamma >= 5% (below: -7.1 dB irrecoverable loss)",
    "",
    "  Gate 2 --> Minimum carrier budget",
    "    What's the minimum photon flux / SNR for the solver to work?",
    "    Example: MRI needs noise sigma < 1% (below: -11.7 dB, correction fails)",
    "",
    "  Gate 3 --> Maximum tolerable calibration error",
    "    How precisely must you calibrate for < 1 dB degradation?",
    "    Example: CT center-of-rotation must be within ~2 px",
    "    Example: CASSI mask alignment must be within < 0.25 px",
    "",
    "These convert 'is this design adequate?' into 3 quantitative pass/fail tests.",
    "All derived from the primitive DAG --- no hardware needed!",
], sz=18, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 28: Design by composition
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Design by Primitive Composition: A Worked Example")
bullets(s, 0.8, 1.3, 11.7, 5.8, [
    "Task: Design a snapshot compressive video imager (like CACTI).",
    "",
    "Step 1: Compose primitives from the library:",
    "  Select:  Modulate M (binary coded mask)  +  Accumulate Sigma (temporal sum)  +  Detect D",
    "  DAG:     M -> Sigma -> D    (3 nodes, depth 3)",
    "",
    "Step 2: Evaluate Gate 1 (Recoverability):",
    "  At compression B=8, measurement matrix retains sufficient rank for ~27 dB (GAP-TV)",
    "",
    "Step 3: Evaluate Gate 2 (Carrier Budget):",
    "  Minimum ~5,000 peak photons per compressed frame for solver to work above noise floor",
    "",
    "Step 4: Evaluate Gate 3 (Calibration Tolerance):",
    "  Mask alignment is dominant sensitivity: sub-pixel shift causes 13 dB degradation!",
    "  Design spec: mask fabrication tolerance < 0.25 px",
    "",
    "Output: Complete design specification BEFORE building anything!",
    "  - Min compression ratio, min photon budget, max alignment error",
    "  - All derived from the primitive DAG + gate analysis",
], sz=17, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 29: Systematic design exploration
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Systematic Design Exploration: 170-Modality Registry")
bullets(s, 0.8, 1.3, 5.8, 5.8, [
    "The registry enables purpose-conditioned",
    "system selection:",
    "",
    "Input: task requirements",
    "  - Target resolution, bandwidth",
    "  - Budget, sample constraints",
    "  - Spectral range, frame rate...",
    "",
    "Output: ranked feasible architectures",
    "  Each with 3-gate design specs",
    "",
    "Example: spectral imager for",
    "precision agriculture",
    "  - 400-1000 nm, >= 20 bands,",
    "    >= 30 fps, field-portable",
], sz=17, color=DARK)
bullets(s, 7.0, 1.3, 5.8, 5.8, [
    "Registry returns candidates:",
    "",
    "  1. CASSI (coded aperture)",
    "     G1: min spatial*spectral density",
    "     G2: min illumination for 30fps",
    "     G3: mask alignment < 0.25 px",
    "",
    "  2. Mosaic filter sensor",
    "     G3: filter registration accuracy",
    "",
    "  3. Filter wheel system",
    "     G1: temporal sampling limit",
    "",
    "Gate-based specs enable quantitative",
    "comparison across architecturally",
    "distinct systems on common footing!",
    "",
    "From TASK -> HARDWARE TOLERANCES",
], sz=17, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 30: Implications for design
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Why This Changes How We Design Imaging Systems")
bullets(s, 0.8, 1.4, 11.7, 5.5, [
    "Before: imaging system design is an open-ended search.",
    "  Intuition-driven, modality-specific, no cross-domain comparison.",
    "",
    "After: imaging system design is a COMBINATORIAL OPTIMIZATION over 11 typed primitives.",
    "  Select primitives, connect DAG, evaluate 3 gates --> done.",
    "",
    "Directly relevant to:",
    "  - Multi-scale camera architectures (Brady et al. 2012):",
    "    Dozens of sub-apertures must be jointly calibrated --> Gate 3 tolerance analysis",
    "  - Snapshot compressive imaging systems (Yuan et al. 2021):",
    "    Measurement operator encodes the entire recovery problem --> Gate 1 + Gate 3",
    "",
    "For AI Scientists: not open-ended function approximation, but a FINITE search",
    "  over 11 known primitives. Identify which subset is active + estimate parameters.",
    "",
    "Broader question: do other measurement domains (seismology, radio astronomy,",
    "particle physics) admit similarly compact grammars?",
], sz=18, color=DARK)
snum(s)

# ══════════════════════════════════════════════════════════════════
# 31: Summary
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s)
title_bar(s, "Summary: The Universal Grammar of Computational Imaging")

# Three pillar summary
box(s, 0.5, 1.4, 3.8, 2.3, RGBColor(0xE0,0xEE,0xFA), GATE1)
tb(s, 0.7, 1.45, 3.4, 0.4, "DIAGNOSE", sz=22, bold=True, color=GATE1, align=PP_ALIGN.CENTER)
bullets(s, 0.7, 1.9, 3.4, 1.6, [
    "Triad diagnoses root cause",
    "Gate 3 dominant 14/14",
    "Cross-over thresholds",
    "8/12 predictions confirmed",
], sz=14, color=DARK, sp=Pt(2))

box(s, 4.8, 1.4, 3.8, 2.3, RGBColor(0xFA,0xE0,0xE0), GATE3)
tb(s, 5.0, 1.45, 3.4, 0.4, "CORRECT", sz=22, bold=True, color=GATE3, align=PP_ALIGN.CENTER)
bullets(s, 5.0, 1.9, 3.4, 1.6, [
    "Solver-agnostic correction",
    "+0.8 to +13.9 dB recovery",
    "12 modalities, 5 carriers",
    "82-100% of oracle ceiling",
], sz=14, color=DARK, sp=Pt(2))

box(s, 9.1, 1.4, 3.8, 2.3, RGBColor(0xE0,0xFA,0xE0), GREEN)
tb(s, 9.3, 1.45, 3.4, 0.4, "DESIGN NEW", sz=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
bullets(s, 9.3, 1.9, 3.4, 1.6, [
    "Compose from 11 primitives",
    "3 gates = design specs",
    "Evaluate before fabrication",
    "170-modality registry",
], sz=14, color=DARK, sp=Pt(2))

# Key numbers
box(s, 0.5, 4.0, 12.4, 3.0, RGBColor(0xF0,0xF0,0xF8), BLUE)
bullets(s, 0.8, 4.1, 11.8, 2.8, [
    "11 primitives --- sufficient and minimal basis for every imaging forward model",
    "3 gates --- exhaustive failure decomposition: MSE = MSE_G1 + MSE_G2 + MSE_G3",
    "170 modality templates across 5 carrier families (+ particle beams)",
    "12 modalities with full 4-scenario validation (2 Nobel Prize techniques, 3 clinical instruments)",
    "+0.8 to +13.9 dB recovery on deployed instruments through operator correction alone",
    "Solver-agnostic: one calibration benefits ALL downstream solvers",
    "NEW: compositional design language turns system design into constrained DAG optimization",
    "",
    "The first universal grammar for validating, correcting, and designing computational imaging systems.",
], sz=17, color=DARK, sp=Pt(3))
snum(s)

# ══════════════════════════════════════════════════════════════════
# 32: Thank you
# ══════════════════════════════════════════════════════════════════
sn += 1; s = new_slide(); add_bg(s, BG_DARK)
tb(s, 0.8, 1.8, 11.7, 1.0, "Thank You!", sz=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 0.8, 2.9, 11.7, 0.8, "Questions & Discussion", sz=32, color=GATE1, align=PP_ALIGN.CENTER)

# Three pillars again
box(s, 1.5, 4.0, 3.0, 0.7, RGBColor(0x1E,0x3A,0x5F), GATE1)
tb(s, 1.5, 4.05, 3.0, 0.6, "DIAGNOSE", sz=22, bold=True, color=GATE1, align=PP_ALIGN.CENTER)
box(s, 5.2, 4.0, 3.0, 0.7, RGBColor(0x3A,0x1E,0x1E), GATE3)
tb(s, 5.2, 4.05, 3.0, 0.6, "CORRECT", sz=22, bold=True, color=GATE3, align=PP_ALIGN.CENTER)
box(s, 8.9, 4.0, 3.0, 0.7, RGBColor(0x1E,0x3A,0x1E), GREEN)
tb(s, 8.9, 4.05, 3.0, 0.6, "DESIGN NEW", sz=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

tb(s, 0.8, 5.2, 11.7, 0.5, "Chengshuai Yang  |  integrityyang@gmail.com", sz=20, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 0.8, 5.7, 11.7, 0.5, "github.com/integritynoble/Physics_World_Model", sz=18, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 0.8, 6.2, 11.7, 0.5, "Acknowledgments: Prof. David Brady (U. Arizona)  |  Prof. Xin Yuan (Westlake)", sz=18, color=GRAY, align=PP_ALIGN.CENTER)
snum(s)

# ── Save ──────────────────────────────────────────────────────────
out_path = OUT / "talk_pwm_flagship_new.pptx"
for suffix in ["_new", "", "_v2", "_v3"]:
    out_path = OUT / f"talk_pwm_flagship{suffix}.pptx"
    try:
        prs.save(str(out_path))
        break
    except PermissionError:
        print(f"(Locked: {out_path.name})")
        continue
print(f"\nPresentation saved: {out_path}")
print(f"Total slides: {sn}")
